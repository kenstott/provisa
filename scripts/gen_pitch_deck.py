# Copyright (c) 2026 Kenneth Stott
# Canary: 9070ddf5-1de2-4070-8190-c78d50f007bc
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""
Generates provisa-overview.pptx — tech exec / architect overview deck.
Usage: python3 scripts/gen_pitch_deck.py
Output: docs/provisa-overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x1E, 0x40)   # deep navy — primary bg
TEAL      = RGBColor(0x00, 0xA8, 0xB5)   # teal accent
TEAL_DARK = RGBColor(0x00, 0x7A, 0x85)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF0, 0xF4, 0xF8)
SLATE     = RGBColor(0x2C, 0x3E, 0x50)
MID_GRAY  = RGBColor(0x8A, 0x9B, 0xB0)
LIGHT_GRAY= RGBColor(0xD0, 0xD8, 0xE4)
DARK_TEXT = RGBColor(0x1A, 0x2A, 0x40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=13, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def navy_bg(slide):
    bg = add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
    return bg


def teal_bar(slide, y=0.55, h=0.06):
    add_rect(slide, 0, y, 13.333, h, fill_color=TEAL)


def slide_header(slide, title, subtitle=None):
    """Standard content slide header: navy top band + title."""
    add_rect(slide, 0, 0, 13.333, 1.15, fill_color=NAVY)
    teal_bar(slide, y=1.15, h=0.05)
    add_textbox(slide, 0.4, 0.15, 12.5, 0.85, title,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.4, 0.75, 12.5, 0.45, subtitle,
                    font_size=14, bold=False, color=TEAL, align=PP_ALIGN.LEFT)


def capability_card(slide, left, top, width, height, title, bullets):
    """Teal-bordered card with title + bullet list."""
    # Card bg
    card = add_rect(slide, left, top, width, height, fill_color=OFFWHITE,
                    line_color=TEAL, line_width=1.5)
    # Title strip
    add_rect(slide, left, top, width, 0.38, fill_color=TEAL)
    add_textbox(slide, left + 0.1, top + 0.04, width - 0.15, 0.32,
                title, font_size=13, bold=True, color=WHITE)
    # Bullets
    txb = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.45),
        Inches(width - 0.22), Inches(height - 0.55)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(11.5)
        run.font.color.rgb = DARK_TEXT
    return card


def arch_box(slide, left, top, width, height, label, sublabel=None,
             fill=TEAL, text_color=WHITE, font_size=13):
    add_rect(slide, left, top, width, height, fill_color=fill,
             line_color=TEAL_DARK, line_width=1)
    y_label = top + (height / 2) - (0.22 if sublabel else 0.15)
    add_textbox(slide, left + 0.05, y_label, width - 0.1, 0.32,
                label, font_size=font_size, bold=True,
                color=text_color, align=PP_ALIGN.CENTER)
    if sublabel:
        add_textbox(slide, left + 0.05, y_label + 0.27, width - 0.1, 0.28,
                    sublabel, font_size=10, bold=False,
                    color=text_color, align=PP_ALIGN.CENTER)


def arrow_down(slide, cx, y, length=0.25, color=TEAL):
    """Simple downward arrow as a thin rectangle + triangle tip."""
    # shaft
    add_rect(slide, cx - 0.02, y, 0.04, length, fill_color=color)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    navy_bg(slide)

    # Left accent bar
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=TEAL)

    # Product name
    add_textbox(slide, 0.5, 1.6, 8.5, 1.4, "Provisa",
                font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tagline
    add_textbox(slide, 0.5, 3.05, 9.0, 0.6,
                "Governed Data Provisioning Platform",
                font_size=26, bold=False, color=TEAL, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(slide, 0.5, 3.75, 8.5, 0.045, fill_color=TEAL)

    # Sub-tagline
    add_textbox(slide, 0.5, 3.9, 10.5, 0.55,
                "Federated query  ·  Semantic model  ·  Pre-approval security  ·  Multi-protocol delivery",
                font_size=15, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Audience label
    add_textbox(slide, 0.5, 6.6, 6, 0.55,
                "Technical Executive & Architect Overview",
                font_size=13, italic=True, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Right graphic: stacked labels
    labels = ["Federation", "Semantic Layer", "Governance", "Delivery"]
    colors = [TEAL, TEAL_DARK, RGBColor(0x0A,0x50,0x70), RGBColor(0x05,0x35,0x55)]
    for i, (lbl, clr) in enumerate(zip(labels, colors)):
        add_rect(slide, 10.2, 1.5 + i * 1.1, 2.8, 0.9, fill_color=clr)
        add_textbox(slide, 10.2, 1.52 + i * 1.1, 2.8, 0.8,
                    lbl, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_product_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "Product Summary", "Core capabilities")

    cards = [
        (
            "Federated Query Engine",
            [
                "Query across any combination of data sources in a single request",
                "Distributed execution via federation engine — no single-process ceiling",
                "Predicate pushdown to source systems; cost-based join strategy selection",
                "Memory management with spill-to-disk for large analytical workloads",
                "Single-source queries bypass federation — sub-100ms RDBMS latency",
                "Cross-source queries: 300–500ms typical; large results async to blob storage",
            ]
        ),
        (
            "Semantic Layer Management",
            [
                "Register sources, tables, and relationships through a governed UI",
                "GraphQL schema auto-generated from registration model — never hand-authored",
                "Cross-source relationships defined semantically, not inferred from FK constraints",
                "Schema reflects business intent; stewards control what the surface exposes",
                "NoSQL sources automatically materialized to Parquet for unified querying",
                "Schema changes flag dependent relationships and registry entries for re-review",
            ]
        ),
        (
            "Governance & Security",
            [
                "Pre-approval model: production queries are pre-authorized, not runtime-evaluated",
                "Persisted query registry — every production query is enumerable and auditable",
                "Three independent enforcement layers: pre-approval, schema visibility, SQL enforcement",
                "Row-level security injected at executor; column security at schema generation",
                "Authorization failures happen at development time, not production time",
                "User rights and query governance are orthogonal — neither substitutes for the other",
            ]
        ),
        (
            "Query Languages & Delivery",
            [
                "GraphQL: primary interface for application consumers and UI backends",
                "gRPC / Arrow Flight: zero-copy columnar streaming for pipelines and analytics",
                "JDBC driver: direct connection from BI tools (Tableau, Power BI, Looker)",
                "Presigned URL redirect: large bulk exports to blob storage (JSON, NDJSON, Parquet, CSV)",
                "Cypher: graph query language for relationship traversal use cases",
                "All entry points share the same governed execution pipeline",
            ]
        ),
        (
            "Breadth of Use Cases",
            [
                "Operational: real-time API serving for web and mobile applications",
                "Analytical: enterprise-scale cross-source reporting and aggregations",
                "Data provisioning: governed bulk export for downstream consumers and pipelines",
                "BI enablement: single governed layer eliminating inconsistent report definitions",
                "Self-service: developers compose precise queries without involving data teams",
                "Open source core — enterprise-class replacement for Hasura v2",
            ]
        ),
        (
            "Persisted Query Registry",
            [
                "Queries submitted for steward review before reaching production",
                "Approved queries receive a stable identifier — query text never sent in production",
                "Steward approves: table scope, parameter schema, permitted output types",
                "Audit trail: who defined, who approved, when, against which schema version",
                "Deprecated queries return clear errors directing clients to replacements",
                "Clients can restrict within approved ceiling; cannot exceed it",
            ]
        ),
    ]

    cols = 3
    card_w = 4.1
    card_h = 2.55
    gap_x = 0.15
    gap_y = 0.18
    start_x = 0.22
    start_y = 1.28

    for i, (title, bullets) in enumerate(cards):
        col = i % cols
        row = i // cols
        lx = start_x + col * (card_w + gap_x)
        ty = start_y + row * (card_h + gap_y)
        capability_card(slide, lx, ty, card_w, card_h, title, bullets)


def slide_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "Deployment Options", "From laptop to enterprise cluster")

    options = [
        {
            "title": "Single Node",
            "icon": "[ ]",
            "tag": "Development / Small Prod",
            "points": [
                "Docker Compose (core stack)",
                "All services in one compose file: federation engine, API, UI, storage",
                "AppImage packaging for air-gapped / on-prem bare-metal",
                "Single command startup; health-checked service ordering",
                "Suitable for dev, demo, and small production workloads",
            ]
        },
        {
            "title": "Multi-Node — Cloud (Terraform)",
            "icon": "[ ][ ]",
            "tag": "AWS · Azure · GCP",
            "points": [
                "1 primary node (coordinator + API + UI) + N worker nodes",
                "Workers are memory-optimized instances (128–256 GB RAM recommended)",
                "Separate instance types for primary (general) vs workers (memory-optimized)",
                "External load balancer provisioned for API (8000) and Arrow Flight (8815)",
                "Interactive deploy.sh wizard: region, node count, instance sizing",
                "AppImage pulled from cloud object storage on startup — no container registry",
            ]
        },
        {
            "title": "Kubernetes / Helm",
            "icon": "⎈",
            "tag": "Cloud-native orchestration",
            "points": [
                "Helm chart for full stack deployment",
                "Horizontal pod autoscaling for worker tier",
                "Federation engine workers: pull-based discovery — new pods join automatically",
                "Proactive scaling recommended over reactive (cold start ~60s per worker)",
                "Resource requests/limits aligned to memory-optimized profiles",
                "Ingress for GraphQL API; ClusterIP for internal Arrow Flight",
            ]
        },
        {
            "title": "On-Premises / Hypervisor",
            "icon": "⬚",
            "tag": "Enterprise private cloud",
            "points": [
                "AppImage runs on any Linux VM — no container runtime required",
                "Role flags: --role primary | --role secondary --primary-ip <ip>",
                "Workers self-register with coordinator on startup",
                "Terraform modules adaptable to on-prem providers (vSphere, OpenStack)",
                "No cost when idle on private hypervisor — VM resources shared with host",
                "Recommended: start with one beefy node; add workers when concurrency demands it",
            ]
        },
    ]

    card_w = 6.2
    card_h = 3.65
    gap = 0.25
    positions = [
        (0.2, 1.3),
        (6.65, 1.3),
        (0.2, 5.1 - 0.15),
        (6.65, 5.1 - 0.15),
    ]

    # Only 2 rows visible — adjust y
    positions = [
        (0.2, 1.28),
        (6.7, 1.28),
        (0.2, 4.98),
        (6.7, 4.98),
    ]

    for opt, (lx, ty) in zip(options, positions):
        # Card bg
        add_rect(slide, lx, ty, card_w, card_h, fill_color=WHITE,
                 line_color=TEAL, line_width=1.5)
        # Header strip
        add_rect(slide, lx, ty, card_w, 0.48, fill_color=NAVY)
        add_textbox(slide, lx + 0.12, ty + 0.05, card_w - 0.55, 0.35,
                    opt["title"], font_size=14, bold=True, color=WHITE)
        add_textbox(slide, lx + card_w - 2.1, ty + 0.07, 2.0, 0.32,
                    opt["tag"], font_size=10, italic=True, color=TEAL,
                    align=PP_ALIGN.RIGHT)

        # Bullet list
        txb = slide.shapes.add_textbox(
            Inches(lx + 0.15), Inches(ty + 0.54),
            Inches(card_w - 0.25), Inches(card_h - 0.62)
        )
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for pt in opt["points"]:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"• {pt}"
            run.font.size = Pt(11.5)
            run.font.color.rgb = DARK_TEXT


def slide_arch_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "Architecture Overview", "Four-layer pipeline — each layer has a single responsibility")

    # ── Layer stack (left 8.5 inches wide, centered) ──────────────────────────
    layers = [
        ("Registration & Governance Layer",
         "Source onboarding · Table/relationship registration · Column & row security · Persisted query registry",
         RGBColor(0x0F, 0x3A, 0x6E)),
        ("GraphQL Compiler",
         "Schema generation from registration model · GraphQL → PG-style SQL · Two-pass: SDL gen + query compilation",
         RGBColor(0x09, 0x5A, 0x7A)),
        ("Transpilation & Routing (SQLGlot)",
         "PG-SQL → Trino SQL or target RDBMS dialect · Single-source vs cross-source routing decision",
         RGBColor(0x00, 0x7A, 0x85)),
        ("Execution Layer",
         "Federation engine (cross-source) · Direct RDBMS (single-source, mutations) · Blob storage redirect (large results)",
         RGBColor(0x00, 0x9A, 0x8A)),
    ]

    box_w = 8.3
    box_h = 0.88
    lx = 0.2
    start_y = 1.28
    gap = 0.14

    for i, (title, sub, clr) in enumerate(layers):
        ty = start_y + i * (box_h + gap)
        add_rect(slide, lx, ty, box_w, box_h, fill_color=clr)
        add_textbox(slide, lx + 0.15, ty + 0.07, box_w - 0.25, 0.35,
                    title, font_size=13, bold=True, color=WHITE)
        add_textbox(slide, lx + 0.15, ty + 0.45, box_w - 0.25, 0.38,
                    sub, font_size=10.5, bold=False, color=LIGHT_GRAY)
        # Arrow between layers
        if i < len(layers) - 1:
            arrow_y = ty + box_h
            add_rect(slide, lx + box_w / 2 - 0.025, arrow_y, 0.05, gap - 0.02,
                     fill_color=TEAL)

    # ── Entry points (right column) ───────────────────────────────────────────
    add_textbox(slide, 9.0, 1.2, 4.1, 0.35,
                "CLIENT ENTRY POINTS", font_size=11, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    entry_points = [
        ("GraphQL API", "Application consumers\nUI backends · developers\nport 8000"),
        ("Arrow Flight / gRPC", "Pipelines · analytics runtimes\nZero-copy columnar stream\nport 8815"),
        ("JDBC Driver", "BI tools — Tableau, Power BI\nLooker, Excel\nTrino-compatible"),
        ("Presigned URL", "Bulk export to blob storage\nJSON · NDJSON · Parquet · CSV\nAsync large results"),
    ]

    ep_w = 3.95
    ep_h = 1.1
    ep_x = 9.1
    ep_start_y = 1.55
    ep_gap = 0.18

    for i, (ep_title, ep_sub) in enumerate(entry_points):
        ty = ep_start_y + i * (ep_h + ep_gap)
        add_rect(slide, ep_x, ty, ep_w, ep_h, fill_color=WHITE,
                 line_color=TEAL, line_width=1.2)
        add_rect(slide, ep_x, ty, ep_w, 0.3, fill_color=TEAL)
        add_textbox(slide, ep_x + 0.1, ty + 0.02, ep_w - 0.12, 0.28,
                    ep_title, font_size=11, bold=True, color=WHITE)
        add_textbox(slide, ep_x + 0.1, ty + 0.33, ep_w - 0.15, 0.72,
                    ep_sub, font_size=10, color=DARK_TEXT)

    # ── Data source badges (bottom) ───────────────────────────────────────────
    add_textbox(slide, 0.2, 6.25, 8.5, 0.3,
                "DATA SOURCES", font_size=10, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    sources = ["PostgreSQL", "MySQL / SQL Server", "Apache Iceberg / S3",
               "MongoDB", "Kafka", "Elasticsearch", "Snowflake / BigQuery", "+ any Trino connector"]
    src_w = 1.0
    src_h = 0.42
    src_gap = 0.025
    total_w = len(sources) * src_w + (len(sources) - 1) * src_gap
    sx = (8.7 - total_w) / 2
    for j, src in enumerate(sources):
        bx = sx + j * (src_w + src_gap)
        clr = TEAL_DARK if j < 7 else NAVY
        add_rect(slide, bx, 6.58, src_w, src_h, fill_color=clr)
        add_textbox(slide, bx + 0.02, 6.6, src_w - 0.04, src_h - 0.06,
                    src, font_size=8.5, bold=False, color=WHITE,
                    align=PP_ALIGN.CENTER)


def slide_query_paths(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "Query Execution Paths", "Structural routing — deterministic, not configurable")

    paths = [
        {
            "title": "Direct Path — Single Source",
            "color": RGBColor(0x09, 0x5A, 0x7A),
            "latency": "< 100ms",
            "x": 0.2,
            "steps": [
                "Client submits pre-approved query ID + bindings",
                "Executor validates registry membership",
                "Compiler produces PG-style SQL",
                "SQLGlot transpiles to target RDBMS dialect",
                "Executes directly against source via warm connection pool",
                "Inline result returned to client",
            ]
        },
        {
            "title": "Federation Path — Cross Source",
            "color": RGBColor(0x00, 0x7A, 0x85),
            "latency": "300–500ms",
            "x": 4.6,
            "steps": [
                "Client submits pre-approved query ID + bindings",
                "Executor validates registry membership",
                "Compiler produces PG-style SQL",
                "SQLGlot transpiles to federation engine SQL",
                "Federation engine: distributed cross-source execution",
                "Inline or redirected result returned",
            ]
        },
        {
            "title": "Large Result Path — Blob Redirect",
            "color": RGBColor(0x00, 0x9A, 0x8A),
            "latency": "Seconds–minutes (async)",
            "x": 9.0,
            "steps": [
                "Client requests pre-approved query with redirect preference",
                "Executor validates registry + output type approval",
                "Federation engine executes, writes to blob storage",
                "Presigned URL returned (TTL-bounded)",
                "Client downloads directly from storage",
                "Formats: JSON · NDJSON · Parquet · CSV",
            ]
        },
    ]

    card_w = 4.0
    for path in paths:
        lx = path["x"]
        ty = 1.28

        # Header
        add_rect(slide, lx, ty, card_w, 0.52, fill_color=path["color"])
        add_textbox(slide, lx + 0.1, ty + 0.04, card_w - 0.55, 0.38,
                    path["title"], font_size=13, bold=True, color=WHITE)
        add_textbox(slide, lx + card_w - 1.55, ty + 0.08, 1.45, 0.35,
                    path["latency"], font_size=10, bold=True,
                    color=TEAL, align=PP_ALIGN.RIGHT)

        # Steps
        step_h = 0.72
        for i, step in enumerate(path["steps"]):
            sy = ty + 0.58 + i * (step_h + 0.04)
            # unused — kept for future shading
            _ = i
            add_rect(slide, lx + 0.05, sy, card_w - 0.1, step_h,
                     fill_color=WHITE, line_color=LIGHT_GRAY, line_width=0.5)
            # Step number
            add_rect(slide, lx + 0.05, sy, 0.32, step_h, fill_color=path["color"])
            add_textbox(slide, lx + 0.05, sy + 0.18, 0.32, 0.35,
                        str(i + 1), font_size=13, bold=True,
                        color=WHITE, align=PP_ALIGN.CENTER)
            add_textbox(slide, lx + 0.42, sy + 0.1, card_w - 0.55, step_h - 0.15,
                        step, font_size=10.5, color=DARK_TEXT)

    # Security callout bar
    add_rect(slide, 0.2, 6.7, 12.9, 0.55, fill_color=NAVY)
    add_textbox(slide, 0.35, 6.75, 12.6, 0.42,
                "All paths: pre-approval check · schema visibility enforcement · RLS injection · column security — same pipeline regardless of entry point or routing",
                font_size=11.5, bold=False, color=TEAL, align=PP_ALIGN.CENTER)


def slide_security_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "Security Model", "Three independent enforcement layers — no single layer depends on the others")

    layers = [
        {
            "num": "1",
            "title": "Pre-Approval Layer",
            "desc": (
                "No query executes in production without explicit steward authorization. "
                "The production endpoint accepts only registry identifiers — not query text. "
                "Every executable operation is enumerable and auditable before a single "
                "production request is made. Cannot be bypassed by privilege escalation."
            ),
            "threat": "Eliminates: arbitrary query injection, privilege-escalation queries, unknown attack surface",
        },
        {
            "num": "2",
            "title": "Schema Visibility Layer",
            "desc": (
                "The GraphQL SDL is generated per-role from column visibility rules in the "
                "registration model. Unauthorized tables and columns do not exist in the schema. "
                "A user cannot reference what they cannot see. Rejections happen at compile time — "
                "the invalid reference never reaches the execution backend."
            ),
            "threat": "Eliminates: data enumeration, column inference attacks, unauthorized table access",
        },
        {
            "num": "3",
            "title": "SQL Enforcement Layer",
            "desc": (
                "Row-level security WHERE clauses and column projections are injected at the "
                "executor — after compile, before execution. Applied to every query regardless "
                "of schema state. Ensures no pre-approved query invoked by a narrow-rights user "
                "can exceed that user's authorized data boundary."
            ),
            "threat": "Eliminates: cross-user data leakage, RLS bypass via pre-approved queries",
        },
    ]

    lx = 0.25
    card_w = 12.8
    card_h = 1.62
    start_y = 1.28
    gap = 0.18

    colors = [
        RGBColor(0x0F, 0x3A, 0x6E),
        RGBColor(0x09, 0x5A, 0x7A),
        RGBColor(0x00, 0x7A, 0x85),
    ]

    for i, (layer, clr) in enumerate(zip(layers, colors)):
        ty = start_y + i * (card_h + gap)
        add_rect(slide, lx, ty, card_w, card_h, fill_color=WHITE,
                 line_color=clr, line_width=1.5)
        # Number badge
        add_rect(slide, lx, ty, 0.55, card_h, fill_color=clr)
        add_textbox(slide, lx, ty + card_h / 2 - 0.22, 0.55, 0.44,
                    layer["num"], font_size=22, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, lx + 0.65, ty + 0.08, 11.8, 0.36,
                    layer["title"], font_size=14, bold=True, color=clr)
        # Description
        add_textbox(slide, lx + 0.65, ty + 0.45, 11.8, 0.72,
                    layer["desc"], font_size=11, color=DARK_TEXT)
        # Threat tag
        add_textbox(slide, lx + 0.65, ty + 1.2, 11.8, 0.35,
                    layer["threat"], font_size=10, italic=True,
                    color=TEAL_DARK)

    # Orthogonality callout
    add_rect(slide, 0.25, 6.42, 12.8, 0.72, fill_color=NAVY)
    add_textbox(slide, 0.45, 6.5, 12.4, 0.55,
                "User rights govern what data flows through an approved operation.  "
                "Query governance governs which operations exist.  "
                "The two layers are orthogonal — neither substitutes for the other.",
                font_size=12, bold=False, color=TEAL, align=PP_ALIGN.CENTER)


def slide_devops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
    slide_header(slide, "DevOps & Operations", "Deployment lifecycle · Monitoring · Scaling")

    # ── Left: Deployment Lifecycle ────────────────────────────────────────────
    add_rect(slide, 0.2, 1.28, 6.0, 0.38, fill_color=NAVY)
    add_textbox(slide, 0.3, 1.32, 5.8, 0.3,
                "DEPLOYMENT LIFECYCLE", font_size=11, bold=True, color=WHITE)

    lifecycle = [
        ("Build", "CI produces AppImage artifact; pushed to cloud object storage (S3 / GCS / Azure Blob)"),
        ("Provision", "Terraform apply — VMs provisioned, LB configured, firewall rules set"),
        ("Bootstrap", "Startup script: pull AppImage from storage, chmod +x, launch with --role flag"),
        ("Scale Out", "Add worker nodes: Terraform node_count + 1; workers self-register with coordinator"),
        ("Upgrade", "Rolling: terminate workers one at a time, pull new AppImage, restart; coordinator last"),
        ("Rollback", "Re-point object storage path to prior AppImage version; restart nodes"),
    ]

    for i, (step, desc) in enumerate(lifecycle):
        ty = 1.72 + i * 0.72
        clr = TEAL if i % 2 == 0 else TEAL_DARK
        add_rect(slide, 0.2, ty, 1.1, 0.6, fill_color=clr)
        add_textbox(slide, 0.22, ty + 0.12, 1.06, 0.38,
                    step, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.38, ty + 0.06, 4.7, 0.54,
                    desc, font_size=10.5, color=DARK_TEXT)

    # ── Right: Monitoring Stack ───────────────────────────────────────────────
    add_rect(slide, 6.7, 1.28, 6.4, 0.38, fill_color=NAVY)
    add_textbox(slide, 6.82, 1.32, 6.2, 0.3,
                "OBSERVABILITY STACK", font_size=11, bold=True, color=WHITE)

    monitoring = [
        (
            "Metrics — Prometheus + Grafana",
            [
                "API: request rate, latency p50/p95/p99, error rate, active connections",
                "GraphQL: query duration p95",
                "DB: client operation duration p95",
                "Federation Engine: running queries, queued queries, cluster memory free, blocked nodes",
                "OTel Collector: spans accepted/refused",
            ]
        ),
        (
            "Distributed Tracing — Tempo",
            [
                "End-to-end traces: client request → compiler → SQLGlot → execution backend",
                "Federation engine instrumented via OTel Java agent",
                "Trace correlation: service.name=federation-engine in all spans",
            ]
        ),
        (
            "Alerting Guidance",
            [
                "Queued queries > 0 sustained → concurrency bottleneck → add worker nodes",
                "Cluster memory free dropping → memory pressure → vertical scale or add workers",
                "Blocked nodes > 0 → memory crisis → scale immediately",
                "Error rate spike → check registry / compiler logs",
            ]
        ),
    ]

    my = 1.72
    for (title, points) in monitoring:
        # Section header
        add_rect(slide, 6.7, my, 6.4, 0.32, fill_color=TEAL_DARK)
        add_textbox(slide, 6.82, my + 0.03, 6.2, 0.27,
                    title, font_size=11, bold=True, color=WHITE)
        my += 0.34

        txb = slide.shapes.add_textbox(
            Inches(6.82), Inches(my),
            Inches(6.2), Inches(len(points) * 0.3 + 0.05)
        )
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for pt in points:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(1)
            run = p.add_run()
            run.text = f"• {pt}"
            run.font.size = Pt(10.5)
            run.font.color.rgb = DARK_TEXT
        my += len(points) * 0.295 + 0.12

    # Sizing guidance footer
    add_rect(slide, 0.2, 6.72, 12.9, 0.52, fill_color=NAVY)
    add_textbox(slide, 0.35, 6.78, 12.6, 0.38,
                "Sizing guidance: vertical scale first (RAM is the primary constraint) · "
                "Queued queries = concurrency problem → add workers · "
                "Memory free dropping = vertical scale or add workers",
                font_size=11, color=TEAL, align=PP_ALIGN.CENTER)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_bg(slide)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=TEAL)

    add_textbox(slide, 0.5, 1.8, 12.0, 1.2, "Provisa",
                font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, 0.5, 3.1, 10.0, 0.05, fill_color=TEAL)

    pillars = [
        ("Federated", "Query any source, any combination"),
        ("Governed", "Pre-approval model — production is auditable before first request"),
        ("Delivered", "GraphQL · Arrow Flight · JDBC · Bulk export"),
        ("Operated", "Terraform · Kubernetes · Docker · AppImage"),
    ]

    px = 0.5
    for title, sub in pillars:
        add_textbox(slide, px, 3.3, 3.0, 0.4,
                    title, font_size=18, bold=True, color=TEAL)
        add_textbox(slide, px, 3.72, 3.0, 0.45,
                    sub, font_size=11, color=LIGHT_GRAY)
        px += 3.15

    add_textbox(slide, 0.5, 6.4, 12.0, 0.55,
                "Open source · Enterprise-class · Purpose-built governed data provisioning",
                font_size=14, italic=True, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_product_summary(prs)
    slide_deployment(prs)
    slide_arch_overview(prs)
    slide_query_paths(prs)
    slide_security_model(prs)
    slide_devops(prs)
    slide_closing(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "docs", "provisa-overview.pptx")
    out = os.path.normpath(out)
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
